import io
import os
import re
import uuid
import json
import asyncio
import aiofiles
from typing import List, Tuple
import pypdf
from fastapi import UploadFile
from dotenv import load_dotenv

from services.embeddings import embed_texts, embed_query
from services.vectorstore import add_chunks, add_image_chunks, query_collection, delete_collection

load_dotenv()

UPLOAD_DIR      = os.getenv("UPLOAD_DIR", "./uploads")
CHUNK_SIZE      = 400    # tokens (approximated by words)
CHUNK_OVERLAP   = 50
MAX_FILE_MB     = int(os.getenv("MAX_UPLOAD_SIZE_MB", 10))

ALLOWED_EXTENSIONS = {"pdf", "txt", "md", "docx", "pptx", "csv", "xlsx", "xls", "json", "png", "jpg", "jpeg"}


# ── Text extraction ───────────────────────────────────────────────────────────

def extract_text_from_pdf_bytes(pdf_bytes: bytes) -> str:
    """Extract text from raw PDF bytes using pypdf (fallback)."""
    reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
    return "\n\n".join(text_parts)


def extract_text_from_pdf_page_aware(pdf_path: str) -> List[Tuple[str, int]]:
    """
    Extract text from a PDF with page-level awareness using PyMuPDF.

    Returns a list of (page_text, page_num) tuples (1-indexed).
    Falls back to pypdf if PyMuPDF is unavailable.
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf_path)
        pages = []
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                pages.append((text, page_num))
        doc.close()
        return pages
    except ImportError:
        # Fallback: pypdf without page numbers
        with open(pdf_path, "rb") as f:
            text = extract_text_from_pdf_bytes(f.read())
        return [(text, 1)]


def extract_text_from_pdf(path: str) -> str:
    with open(path, "rb") as f:
        return extract_text_from_pdf_bytes(f.read())


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_docx(path: str) -> str:
    """Extract text from a .docx Word document."""
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return "\n\n".join(paragraphs)
    except ImportError:
        raise ValueError("python-docx is not installed. Run: pip install python-docx")


def extract_text_from_pptx(path: str) -> str:
    """Extract text from a .pptx PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"Slide {i}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if len(slide_parts) > 1:
                slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)
    except ImportError:
        raise ValueError("python-pptx is not installed. Run: pip install python-pptx")


def extract_text_from_csv(path: str) -> str:
    """Extract text from a .csv file."""
    try:
        import pandas as pd
        df = pd.read_csv(path)
        return f"Columns: {', '.join(str(c) for c in df.columns)}\n\n{df.to_string(index=False)}"
    except ImportError:
        return extract_text_from_txt(path)


def extract_text_from_excel(path: str) -> str:
    """Extract text from .xlsx / .xls files."""
    try:
        import pandas as pd
        xl = pd.ExcelFile(path)
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            parts.append(f"Sheet: {sheet}\nColumns: {', '.join(str(c) for c in df.columns)}\n{df.to_string(index=False)}")
        return "\n\n".join(parts)
    except ImportError:
        raise ValueError("pandas / openpyxl not installed. Run: pip install pandas openpyxl")


def extract_text_from_json(path: str) -> str:
    """Convert JSON file content to readable text."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        try:
            data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            return f.read()


def extract_text_from_image(path: str) -> str:
    """Extract text/description from an image using Groq's vision model."""
    import base64
    import asyncio
    from services.llm import complete_with_vision

    with open(path, "rb") as f:
        image_bytes = f.read()
    b64 = base64.b64encode(image_bytes).decode("utf-8")

    ext = path.rsplit(".", 1)[-1].lower()
    media_type = "image/png" if ext == "png" else "image/jpeg"

    prompt = (
        "You are an OCR and document analysis assistant. "
        "Extract ALL text visible in this image. If it contains handwritten notes, "
        "equations, diagrams, or charts, describe them in detail. "
        "If there is no text, provide a detailed description of the image content. "
        "Be thorough and accurate."
    )

    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as pool:
                future = pool.submit(asyncio.run, complete_with_vision(b64, prompt, image_media_type=media_type, max_tokens=2000))
                return future.result(timeout=60)
        else:
            return asyncio.run(complete_with_vision(b64, prompt, image_media_type=media_type, max_tokens=2000))
    except Exception as e:
        return f"[Image content - OCR extraction failed: {str(e)}]"


def extract_text(path: str, filename: str) -> str:
    """Route extraction to the correct handler based on file extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"
    if ext == "pdf":
        return extract_text_from_pdf(path)
    elif ext == "docx":
        return extract_text_from_docx(path)
    elif ext == "pptx":
        return extract_text_from_pptx(path)
    elif ext in ("csv",):
        return extract_text_from_csv(path)
    elif ext in ("xlsx", "xls"):
        return extract_text_from_excel(path)
    elif ext == "json":
        return extract_text_from_json(path)
    elif ext in ("png", "jpg", "jpeg"):
        return extract_text_from_image(path)
    else:
        return extract_text_from_txt(path)


# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping word-based chunks."""
    text = re.sub(r"\s+", " ", text).strip()
    words = text.split()

    if not words:
        return []

    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end == len(words):
            break
        start += chunk_size - overlap

    return chunks


def chunk_pages(
    pages: List[Tuple[str, int]],
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> List[Tuple[str, int]]:
    """
    Chunk page-aware text into overlapping word-based chunks.

    Args:
        pages: list of (page_text, page_num) from extract_text_from_pdf_page_aware

    Returns:
        list of (chunk_text, page_num) — each chunk tagged with its source page.
    """
    result = []
    for page_text, page_num in pages:
        page_text = re.sub(r"\s+", " ", page_text).strip()
        words = page_text.split()
        if not words:
            continue
        start = 0
        while start < len(words):
            end = min(start + chunk_size, len(words))
            chunk = " ".join(words[start:end])
            result.append((chunk, page_num))
            if end == len(words):
                break
            start += chunk_size - overlap
    return result


# ── Main pipeline ─────────────────────────────────────────────────────────────

async def ingest_document(file: UploadFile) -> dict:
    """
    Full ingestion pipeline:
    1. Save file to disk
    2. Extract text (page-aware for PDFs, flat for other types)
    3. Chunk text (with page_num metadata for PDFs)
    4. Generate embeddings
    5. Store in ChromaDB

    Returns metadata dict with doc_id, chunk_count, filepath, and image_count.
    """
    # Validate file size
    content = await file.read()
    size_mb = len(content) / (1024 * 1024)
    if size_mb > MAX_FILE_MB:
        raise ValueError(f"File too large ({size_mb:.1f} MB). Max: {MAX_FILE_MB} MB")

    # Validate extension
    ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else "txt"
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '.{ext}'. "
            f"Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
        )

    # Save file
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    doc_id   = str(uuid.uuid4())
    filepath = os.path.join(UPLOAD_DIR, f"{doc_id}.{ext}")

    async with aiofiles.open(filepath, "wb") as f:
        await f.write(content)

    # Extract text — use page-aware extraction for PDFs
    chunk_page_pairs: List[Tuple[str, int]] = []

    if ext == "pdf":
        pages = await asyncio.to_thread(extract_text_from_pdf_page_aware, filepath)
        if not pages:
            raise ValueError("Could not extract text from PDF. Ensure it is not a scanned image.")
        chunk_page_pairs = chunk_pages(pages)
    else:
        raw_text = extract_text(filepath, file.filename)
        if not raw_text.strip():
            raise ValueError("Could not extract text from document. Ensure it is not a scanned image.")
        chunks = chunk_text(raw_text)
        chunk_page_pairs = [(c, 1) for c in chunks]

    if not chunk_page_pairs:
        raise ValueError("Document appears to be empty after text extraction.")

    chunks_only = [c for c, _ in chunk_page_pairs]

    # Embed (wrapped in thread to avoid blocking event loop)
    embeddings = await asyncio.to_thread(embed_texts, chunks_only)

    # Store in Chroma with metadata (page_num included for PDFs)
    metadatas = [
        {"chunk_index": i, "filename": file.filename, "doc_id": doc_id, "page_num": pn}
        for i, (_, pn) in enumerate(chunk_page_pairs)
    ]
    add_chunks(
        collection_name=doc_id,
        chunks=chunks_only,
        embeddings=embeddings,
        metadatas=metadatas,
    )

    # ── Image extraction and indexing ───────────────────────────────────────
    image_count = 0
    if ext == "pdf":
        try:
            # 1. Extract raw images from PDF pages
            raw_images = await asyncio.to_thread(extract_and_store_images, filepath, doc_id)
            
            if raw_images:
                # 2. Classify and generate semantic captions
                classified_images = await classify_images_batch(raw_images)
                
                # Keep only those that successfully got a caption
                valid_images = [img for img in classified_images if img.get("caption")]
                
                if valid_images:
                    captions = [img["caption"] for img in valid_images]
                    
                    # 3. Embed captions
                    image_embeddings = await asyncio.to_thread(embed_texts, captions)
                    
                    # 4. Store in ChromaDB
                    add_image_chunks(
                        collection_name=f"{doc_id}_images",
                        captions=captions,
                        embeddings=image_embeddings,
                        metadatas=valid_images,
                    )
                    image_count = len(valid_images)
        except Exception as e:
            print(f"[ingest_document] Error processing images: {e}")

    return {
        "doc_id":      doc_id,
        "filename":    file.filename,
        "chunk_count": len(chunks_only),
        "image_count": image_count,
        "filepath":    filepath,
        "ext":         ext,
    }


async def query_document(
    collection_name: str,
    question: str,
    top_k: int = 5,
) -> List[Tuple[str, int, float]]:
    """
    Query a stored document collection.

    Returns list of (chunk_text, page_num, cosine_score).
    """
    q_embedding = await asyncio.to_thread(embed_query, question)
    raw_results = query_collection(
        collection_name=collection_name,
        query_embedding=q_embedding,
        top_k=top_k,
    )
    # Extract (text, page_num, score)
    results = []
    for doc, meta, score in raw_results:
        page_num = int(meta.get("page_num", 1))
        results.append((doc, page_num, score))
    return results


async def delete_document(doc_id: str) -> None:
    """Remove document embeddings from ChromaDB and delete the file."""
    delete_collection(doc_id)
    # Also clean up any image sub-collection
    delete_collection(f"{doc_id}_images")
    for ext in ALLOWED_EXTENSIONS:
        path = os.path.join(UPLOAD_DIR, f"{doc_id}.{ext}")
        if os.path.exists(path):
            os.remove(path)
            break
    # Remove extracted images directory
    import shutil
    img_dir = os.path.join(UPLOAD_DIR, "images", doc_id)
    if os.path.exists(img_dir):
        shutil.rmtree(img_dir, ignore_errors=True)